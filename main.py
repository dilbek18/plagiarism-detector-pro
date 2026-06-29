import os
import re
import html
import json
import time
import math
import asyncio
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple
from urllib.parse import urlencode, quote_plus, urlparse
from difflib import SequenceMatcher

import aiohttp
from bs4 import BeautifulSoup
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field

APP_NAME = "Academic Integrity Search Server"
MAX_TEXT_CHARS = int(os.getenv("MAX_TEXT_CHARS", "25000"))
MAX_FETCH_BYTES = int(os.getenv("MAX_FETCH_BYTES", "650000"))
REQUEST_TIMEOUT = int(os.getenv("REQUEST_TIMEOUT", "12"))

# API keys are read only from Render Environment. They are never returned to frontend.
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "").strip()
GOOGLE_CSE_ID = os.getenv("GOOGLE_CSE_ID", "").strip()  # also called CX
BRAVE_API_KEY = os.getenv("BRAVE_API_KEY", "").strip()
SERPER_API_KEY = os.getenv("SERPER_API_KEY", "").strip()
SEARCHAPI_API_KEY = os.getenv("SEARCHAPI_API_KEY", "").strip()

app = FastAPI(title=APP_NAME, version="2.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
)

STATIC_DIR = os.path.join(os.path.dirname(__file__), "static")
if os.path.isdir(STATIC_DIR):
    app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")

class SearchRequest(BaseModel):
    q: str = Field(..., min_length=2, max_length=400)
    limit: int = Field(8, ge=1, le=20)
    academic: bool = False

class CheckRequest(BaseModel):
    text: str = Field(..., min_length=20)
    max_fragments: int = Field(8, ge=1, le=20)
    limit_per_fragment: int = Field(5, ge=1, le=8)
    academic: bool = True

class SearchResult(BaseModel):
    title: str
    url: str
    snippet: str = ""
    engine: str


def clean_text(text: str) -> str:
    text = html.unescape(text or "")
    text = re.sub(r"\s+", " ", text).strip()
    return text[:MAX_TEXT_CHARS]


def normalize_url(url: str) -> str:
    if not url:
        return ""
    url = url.strip()
    if url.startswith("//"):
        url = "https:" + url
    return url


def result_key(item: Dict[str, str]) -> str:
    url = normalize_url(item.get("url", ""))
    parsed = urlparse(url)
    return (parsed.netloc.lower() + parsed.path.rstrip("/")).strip()


def dedupe_results(results: List[Dict[str, str]], limit: int) -> List[Dict[str, str]]:
    seen = set()
    out = []
    for r in results:
        url = normalize_url(r.get("url", ""))
        if not url or not url.startswith(("http://", "https://")):
            continue
        key = result_key({**r, "url": url})
        if not key or key in seen:
            continue
        seen.add(key)
        title = clean_text(r.get("title", ""))[:220] or url
        snippet = clean_text(r.get("snippet", ""))[:500]
        out.append({"title": title, "url": url, "snippet": snippet, "engine": r.get("engine", "web")})
        if len(out) >= limit:
            break
    return out


def academic_query(q: str) -> str:
    q = q.strip()
    # Uses search operators supported by Google/Brave/Serper-style engines.
    return f'{q} (site:.edu OR site:.ac.uk OR site:edu OR filetype:pdf OR site:semanticscholar.org OR site:researchgate.net OR site:core.ac.uk)'


async def http_get_json(session: aiohttp.ClientSession, url: str, **kwargs) -> Optional[Dict[str, Any]]:
    try:
        async with session.get(url, timeout=REQUEST_TIMEOUT, **kwargs) as resp:
            if resp.status >= 400:
                return None
            return await resp.json(content_type=None)
    except Exception:
        return None


async def http_post_json(session: aiohttp.ClientSession, url: str, **kwargs) -> Optional[Dict[str, Any]]:
    try:
        async with session.post(url, timeout=REQUEST_TIMEOUT, **kwargs) as resp:
            if resp.status >= 400:
                return None
            return await resp.json(content_type=None)
    except Exception:
        return None


async def search_google(session: aiohttp.ClientSession, q: str, limit: int) -> List[Dict[str, str]]:
    if not GOOGLE_API_KEY or not GOOGLE_CSE_ID:
        return []
    params = {"key": GOOGLE_API_KEY, "cx": GOOGLE_CSE_ID, "q": q, "num": min(limit, 10)}
    data = await http_get_json(session, "https://www.googleapis.com/customsearch/v1", params=params)
    items = (data or {}).get("items", [])
    results = []
    for it in items:
        results.append({
            "title": it.get("title", ""),
            "url": it.get("link", ""),
            "snippet": it.get("snippet", ""),
            "engine": "Google Programmable Search",
        })
    return results


async def search_brave(session: aiohttp.ClientSession, q: str, limit: int) -> List[Dict[str, str]]:
    if not BRAVE_API_KEY:
        return []
    headers = {"Accept": "application/json", "X-Subscription-Token": BRAVE_API_KEY}
    params = {"q": q, "count": min(limit, 20), "safesearch": "moderate"}
    data = await http_get_json(session, "https://api.search.brave.com/res/v1/web/search", headers=headers, params=params)
    items = ((data or {}).get("web") or {}).get("results", [])
    results = []
    for it in items:
        results.append({
            "title": it.get("title", ""),
            "url": it.get("url", ""),
            "snippet": it.get("description", ""),
            "engine": "Brave Search",
        })
    return results


async def search_serper(session: aiohttp.ClientSession, q: str, limit: int) -> List[Dict[str, str]]:
    if not SERPER_API_KEY:
        return []
    headers = {"X-API-KEY": SERPER_API_KEY, "Content-Type": "application/json"}
    payload = {"q": q, "num": min(limit, 10)}
    data = await http_post_json(session, "https://google.serper.dev/search", headers=headers, json=payload)
    items = (data or {}).get("organic", [])
    results = []
    for it in items:
        results.append({
            "title": it.get("title", ""),
            "url": it.get("link", ""),
            "snippet": it.get("snippet", ""),
            "engine": "Serper Google Results",
        })
    return results


async def search_searchapi_duckduckgo(session: aiohttp.ClientSession, q: str, limit: int) -> List[Dict[str, str]]:
    # Optional paid provider. Useful when Google/Brave keys are absent.
    if not SEARCHAPI_API_KEY:
        return []
    params = {"engine": "duckduckgo", "q": q, "api_key": SEARCHAPI_API_KEY}
    data = await http_get_json(session, "https://www.searchapi.io/api/v1/search", params=params)
    items = (data or {}).get("organic_results", [])
    results = []
    for it in items[:limit]:
        results.append({
            "title": it.get("title", ""),
            "url": it.get("link", ""),
            "snippet": it.get("snippet", ""),
            "engine": "DuckDuckGo via SearchAPI",
        })
    return results


async def search_duckduckgo_lite(session: aiohttp.ClientSession, q: str, limit: int) -> List[Dict[str, str]]:
    # Free fallback. It does not expose API keys. It may be less stable than official/paid APIs.
    url = "https://lite.duckduckgo.com/lite/?" + urlencode({"q": q})
    try:
        async with session.get(url, timeout=REQUEST_TIMEOUT, headers={"User-Agent": "Mozilla/5.0"}) as resp:
            if resp.status >= 400:
                return []
            text = await resp.text(errors="ignore")
    except Exception:
        return []
    soup = BeautifulSoup(text, "lxml")
    results = []
    for a in soup.select("a[href]"):
        href = normalize_url(a.get("href", ""))
        title = clean_text(a.get_text(" "))
        if not title or not href.startswith(("http://", "https://")):
            continue
        if "duckduckgo.com" in urlparse(href).netloc:
            continue
        parent = a.find_parent("td") or a.parent
        snippet = clean_text(parent.get_text(" ") if parent else "")
        results.append({"title": title, "url": href, "snippet": snippet, "engine": "DuckDuckGo Lite fallback"})
        if len(results) >= limit:
            break
    return results


async def multi_search(q: str, limit: int = 8, academic: bool = False) -> Dict[str, Any]:
    q = clean_text(q)[:400]
    if academic:
        q2 = academic_query(q)
    else:
        q2 = q
    timeout = aiohttp.ClientTimeout(total=REQUEST_TIMEOUT + 5)
    connector = aiohttp.TCPConnector(limit=12, ssl=False)
    async with aiohttp.ClientSession(timeout=timeout, connector=connector) as session:
        tasks = [
            search_google(session, q2, limit),
            search_brave(session, q2, limit),
            search_serper(session, q2, limit),
            search_searchapi_duckduckgo(session, q2, limit),
        ]
        # Use DDG Lite fallback only when no API keys are set or other engines give very few results.
        engine_results = await asyncio.gather(*tasks, return_exceptions=True)
        flat: List[Dict[str, str]] = []
        for item in engine_results:
            if isinstance(item, list):
                flat.extend(item)
        if len(flat) < max(3, min(5, limit)):
            flat.extend(await search_duckduckgo_lite(session, q, limit))
    results = dedupe_results(flat, limit)
    enabled_engines = []
    if GOOGLE_API_KEY and GOOGLE_CSE_ID:
        enabled_engines.append("Google")
    if BRAVE_API_KEY:
        enabled_engines.append("Brave")
    if SERPER_API_KEY:
        enabled_engines.append("Serper")
    if SEARCHAPI_API_KEY:
        enabled_engines.append("SearchAPI/DuckDuckGo")
    if not enabled_engines or len(results) < max(3, min(5, limit)):
        enabled_engines.append("DuckDuckGo fallback")
    return {"query": q, "academic_mode": academic, "enabled_engines": enabled_engines, "count": len(results), "results": results}


def split_fragments(text: str, max_fragments: int) -> List[str]:
    text = clean_text(text)
    # Sentence-like splitting first.
    parts = re.split(r"(?<=[.!?。؟])\s+|\n+", text)
    parts = [clean_text(p) for p in parts if len(p.split()) >= 8]
    # Prefer longer fragments but avoid over-long queries.
    parts = sorted(parts, key=lambda s: len(s), reverse=True)
    out = []
    for p in parts:
        p = p[:280]
        if p not in out:
            out.append(p)
        if len(out) >= max_fragments:
            break
    if not out and len(text.split()) >= 8:
        words = text.split()
        for i in range(0, min(len(words), max_fragments * 30), 30):
            frag = " ".join(words[i:i+30])
            if len(frag.split()) >= 8:
                out.append(frag[:280])
    return out[:max_fragments]


def token_set(text: str) -> set:
    return set(re.findall(r"[\w'’\-]+", text.lower(), flags=re.UNICODE))


def similarity_score(a: str, b: str) -> float:
    a = clean_text(a).lower()
    b = clean_text(b).lower()
    if not a or not b:
        return 0.0
    seq = SequenceMatcher(None, a[:1200], b[:1200]).ratio()
    ta, tb = token_set(a), token_set(b)
    jac = len(ta & tb) / max(1, len(ta | tb))
    return round((0.65 * seq + 0.35 * jac) * 100, 2)


async def fetch_page_text(session: aiohttp.ClientSession, url: str) -> str:
    try:
        async with session.get(url, timeout=REQUEST_TIMEOUT, headers={"User-Agent": "Mozilla/5.0"}) as resp:
            if resp.status >= 400:
                return ""
            content = await resp.content.read(MAX_FETCH_BYTES)
            ctype = resp.headers.get("content-type", "")
            if "pdf" in ctype or url.lower().endswith(".pdf"):
                return ""  # Avoid heavy PDF web fetch parsing on free RAM.
            text = content.decode(resp.charset or "utf-8", errors="ignore")
    except Exception:
        return ""
    soup = BeautifulSoup(text, "lxml")
    for tag in soup(["script", "style", "nav", "footer", "header", "noscript"]):
        tag.decompose()
    return clean_text(soup.get_text(" "))


async def check_similarity(text: str, max_fragments: int, limit_per_fragment: int, academic: bool) -> Dict[str, Any]:
    started = time.time()
    text = clean_text(text)
    fragments = split_fragments(text, max_fragments)
    if not fragments:
        return {"overall_similarity": 0, "fragments": [], "sources": [], "message": "Matn juda qisqa yoki fragmentlarga bo‘linmadi."}
    timeout = aiohttp.ClientTimeout(total=REQUEST_TIMEOUT + 15)
    connector = aiohttp.TCPConnector(limit=10, ssl=False)
    all_matches: List[Dict[str, Any]] = []
    async with aiohttp.ClientSession(timeout=timeout, connector=connector) as session:
        for frag in fragments:
            # Query exact-ish phrase and non-quoted version.
            short_q = '"' + frag[:180].replace('"', '') + '"'
            search_data = await multi_search(short_q, limit=limit_per_fragment, academic=academic)
            candidates = search_data.get("results", [])
            if not candidates:
                search_data = await multi_search(frag[:220], limit=limit_per_fragment, academic=academic)
                candidates = search_data.get("results", [])
            for cand in candidates[:limit_per_fragment]:
                snippet_score = similarity_score(frag, cand.get("snippet", "") + " " + cand.get("title", ""))
                page_text = ""
                page_score = 0.0
                # Fetch only if snippet already looks related; saves memory and time.
                if snippet_score >= 18:
                    page_text = await fetch_page_text(session, cand.get("url", ""))
                    page_score = similarity_score(frag, page_text[:1800]) if page_text else 0.0
                score = max(snippet_score, page_score)
                if score >= 22:
                    all_matches.append({
                        "fragment": frag,
                        "source_title": cand.get("title", ""),
                        "source_url": cand.get("url", ""),
                        "engine": cand.get("engine", ""),
                        "similarity": score,
                        "snippet": cand.get("snippet", ""),
                    })
    # Deduplicate matches by URL + fragment.
    unique = []
    seen = set()
    for m in sorted(all_matches, key=lambda x: x["similarity"], reverse=True):
        key = (m["fragment"][:80], result_key({"url": m["source_url"]}))
        if key in seen:
            continue
        seen.add(key)
        unique.append(m)
    fragment_best = {}
    for m in unique:
        fragment_best[m["fragment"]] = max(fragment_best.get(m["fragment"], 0), m["similarity"])
    if fragment_best:
        overall = sum(fragment_best.values()) / len(fragments)
    else:
        overall = 0.0
    overall = round(min(100.0, overall), 2)
    originality = round(max(0, 100 - overall), 2)
    return {
        "overall_similarity": overall,
        "internet_similarity": overall,
        "academic_similarity_estimate": round(overall * 0.65, 2) if academic else 0,
        "originality_score": originality,
        "checked_fragments": len(fragments),
        "matches_found": len(unique),
        "elapsed_seconds": round(time.time() - started, 2),
        "fragments": fragments,
        "sources": unique[:30],
        "note": "Bu natija qidiruv APIlari va ochiq internet snippet/sahifa matniga asoslangan taxminiy o‘xshashlikdir; Turnitin yoki yopiq universitet bazalarini almashtirmaydi.",
    }


def extract_file_text(filename: str, data: bytes) -> str:
    name = (filename or "").lower()
    if len(data) > 7_000_000:
        raise HTTPException(status_code=413, detail="Fayl juda katta. Free Render uchun 7 MB dan kichik fayl yuklang.")
    try:
        if name.endswith(".txt"):
            return data.decode("utf-8", errors="ignore")
        if name.endswith(".html") or name.endswith(".htm"):
            soup = BeautifulSoup(data.decode("utf-8", errors="ignore"), "lxml")
            return soup.get_text(" ")
        if name.endswith(".pdf"):
            from pypdf import PdfReader
            reader = PdfReader(BytesIO(data))
            pages = []
            for page in reader.pages[:25]:
                pages.append(page.extract_text() or "")
            return "\n".join(pages)
        if name.endswith(".docx"):
            from docx import Document
            doc = Document(BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs)
        if name.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(BytesIO(data))
            texts = []
            for slide in prs.slides[:40]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Faylni o‘qib bo‘lmadi: {e}")
    raise HTTPException(status_code=400, detail="Faqat TXT, HTML, PDF, DOCX, PPTX qo‘llanadi.")


@app.get("/", response_class=HTMLResponse)
async def root():
    index = os.path.join(STATIC_DIR, "index.html")
    if os.path.exists(index):
        return FileResponse(index)
    return HTMLResponse(f"<h1>{APP_NAME}</h1><p>Server ishlayapti.</p>")


@app.get("/health")
async def health():
    return {
        "ok": True,
        "app": APP_NAME,
        "engines_configured": {
            "google": bool(GOOGLE_API_KEY and GOOGLE_CSE_ID),
            "brave": bool(BRAVE_API_KEY),
            "serper": bool(SERPER_API_KEY),
            "searchapi_duckduckgo": bool(SEARCHAPI_API_KEY),
            "duckduckgo_fallback": True,
        },
        "api_keys_visible": False,
    }


@app.post("/api/search")
async def api_search(payload: SearchRequest):
    return await multi_search(payload.q, payload.limit, payload.academic)


@app.post("/api/check")
async def api_check(payload: CheckRequest):
    return await check_similarity(payload.text, payload.max_fragments, payload.limit_per_fragment, payload.academic)


@app.post("/api/check-file")
async def api_check_file(
    file: UploadFile = File(...),
    max_fragments: int = Form(8),
    limit_per_fragment: int = Form(5),
    academic: bool = Form(True),
):
    data = await file.read()
    text = clean_text(extract_file_text(file.filename or "", data))
    if len(text.split()) < 15:
        raise HTTPException(status_code=400, detail="Fayldan yetarli matn ajratilmadi.")
    result = await check_similarity(text, max_fragments, limit_per_fragment, academic)
    result["filename"] = file.filename
    result["extracted_chars"] = len(text)
    return result
